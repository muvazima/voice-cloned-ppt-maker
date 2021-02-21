import numpy as np
import pandas as pd
import os
import pickle
import re
from sklearn.datasets import load_files
import glob
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
#import torch
from summarizer import Summarizer
from pptx import Presentation 
from pptx.util import Inches, Pt  

model = Summarizer('distilbert-base-uncased')   

def read_file(filename):

	f = open("/Users/Manam/fyp/Parsed_Papers/"+filename+'.txt', 'r', encoding="utf-8")
	text = str(f.read())
	f.close()
	text = re.sub("\n", " ", text)
	text=re.sub("\t"," ",text)
	text = (text.encode('ascii', 'ignore')).decode("utf-8")
	text= re.sub(' +', ' ', text)
	sections = re.findall(r"@&#\w+@&#", text)
	main_title = re.findall(r'@&#MAIN-TITLE@&#(.*?)@&#', text, flags = re.I)

	d={}
	d['@&#MAIN-TITLE@&#']=re.findall(r'@&#MAIN-TITLE@&#(.*?)@&#', text, flags = re.I)
	for i in sections:
  		d[i]=re.findall(i+'(.*?)@&#', text,  flags = re.I)

	del d['@&#REFERENCES@&#'] 
	return d


def summarize(d):

	donotsummarize=['MAIN-TITLE','HIGHLIGHTS','KEYPHRASES','REFERENCES','ACKNOWLEDGEMENTS']
	lines=[]
	dclean={}
	for i in d:
		iclean=re.sub("@&#","", i)
		if iclean in donotsummarize:
			lines.append(iclean+": "+d[i][0]+"\n\n\n")
			dclean[iclean]=d[i][0]
		else:
			st=model(d[i][0])
			lines.append(iclean+": "+st+"\n\n\n")
			dclean[iclean]=st
	return dclean

def create_ppt(dclean,filename):

	prs = Presentation() 
	first_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(first_slide_layout) 
	slide.shapes.title.text = dclean['MAIN-TITLE']

	for i in dclean: 
  		bullet_slide_layout = prs.slide_layouts[1]
  		slide = prs.slides.add_slide(bullet_slide_layout)
  		shapes = slide.shapes
  		title_shape = shapes.title
  		body_shape = shapes.placeholders[1]
  		title_shape.text = i
  		l=dclean[i].split('.')
  		tf = body_shape.text_frame
  		tf.text = l[0]

  		for j in l[1:len(l)]:
  			p = tf.add_paragraph()
  			p.text =j

	prs.save('/Users/Manam/fyp/PPTs/'+filename+'.pptx')
	#prs.save(filename)

def main():

	print(chr(27) + "[2J")
	print("-------------------------PAPER TO PPT CONVERTER-----------------------------")
	print()
	filename=input('Enter Filename: ')
	d=read_file(filename)
	dclean=summarize(d)
	pptfilename=input('Enter PPT file name: ')
	create_ppt(dclean,pptfilename)
	print("PPT saved.")
	
if __name__ == "__main__":
    main()